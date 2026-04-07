from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("docs/guides/SAP-CashPilot-User-Guide-10-Cas-Usage.pptx")


def set_background(slide, color=(9, 18, 43)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_title(slide, text, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(11.8), Inches(0.9))
        tf_sub = sub_box.text_frame
        tf_sub.clear()
        p_sub = tf_sub.paragraphs[0]
        r_sub = p_sub.add_run()
        r_sub.text = subtitle
        r_sub.font.size = Pt(18)
        r_sub.font.color.rgb = RGBColor(170, 200, 255)
        p_sub.alignment = PP_ALIGN.LEFT


def add_bullets(slide, heading, bullets, x=0.8, y=2.4, w=12.0, h=4.4):
    heading_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.6))
    tf_h = heading_box.text_frame
    tf_h.clear()
    p_h = tf_h.paragraphs[0]
    run_h = p_h.add_run()
    run_h.text = heading
    run_h.font.size = Pt(22)
    run_h.font.bold = True
    run_h.font.color.rgb = RGBColor(255, 214, 102)

    body_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.75), Inches(w), Inches(h))
    tf = body_box.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(232, 240, 255)
        p.space_after = Pt(10)


def add_two_column(slide, left_title, left_points, right_title, right_points):
    add_bullets(slide, left_title, left_points, x=0.7, y=2.0, w=5.8, h=4.6)
    add_bullets(slide, right_title, right_points, x=6.7, y=2.0, w=5.8, h=4.6)


def add_footer(slide, text="CashPilot • Module SAP • Guide Novices"):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(140, 160, 190)
    p.alignment = PP_ALIGN.RIGHT


def add_use_case_slide(slide, numero, titre, contexte, action, resultat, indicateurs):
    add_title(slide, f"Cas d'usage {numero}: {titre}")
    add_two_column(
        slide,
        "Contexte",
        [contexte],
        "Action dans le cockpit SAP",
        [action],
    )
    add_bullets(
        slide,
        "Resultat attendu",
        [resultat, f"Indicateurs: {indicateurs}"],
        x=0.7,
        y=5.25,
        w=11.8,
        h=1.2,
    )
    add_footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s)
    add_title(
        s,
        "SAP dans CashPilot",
        "Presentation + user guide operationnel pour debutants en comptabilite",
    )
    add_bullets(
        s,
        "Ce que vous allez trouver",
        [
            "A quoi sert le module SAP dans CashPilot",
            "Quand l'utiliser dans la vie reelle",
            "10 cas d'usage concrets et actionnables",
            "Parcours accessible pour novices en FR, EN et NL",
            "La plus-value business pour l'entreprise",
        ],
        y=2.6,
    )
    add_footer(s, "CashPilot • SAP Cockpit • Version 2026")

    # Slide 2
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s)
    add_title(s, "Positionnement SAP dans CashPilot")
    add_bullets(
        s,
        "Definition claire",
        [
            "Ce n'est pas un ERP SAP complet",
            "C'est un cockpit de maturite finance/comptabilite",
            "Il structure FI, CO, AA, Consolidation et Close",
            "Il est utilisable en FR, EN et NL dans CashPilot",
            "Il guide les actions via une roadmap mesurable",
        ],
    )
    add_footer(s)

    # Slide 3
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s)
    add_title(s, "Definitions simples des sigles SAP")
    add_bullets(
        s,
        "Glossaire FI / CO / AA",
        [
            "FI (Finance Accounting): comptabilite generale, ecritures, TVA, rapprochements.",
            "CO (Controlling): axes analytiques pour piloter couts, marges et performance.",
            "AA (Asset Accounting): immobilisations, amortissements et suivi des actifs.",
            "Consolidation: pilotage multi-entites et intercompany.",
            "Close: orchestration de la cloture periodique et de la conformite.",
        ],
    )
    add_footer(s)

    # Slide 4
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s)
    add_title(s, "Quand utiliser le cockpit SAP")
    add_bullets(
        s,
        "Moments clefs",
        [
            "Au demarrage d'une nouvelle societe",
            "En revue hebdomadaire ou mensuelle CFO/compta",
            "Avant la cloture de periode",
            "Pendant une transformation finance multi-entites",
            "Pour onboarder des utilisateurs non-comptables",
        ],
    )
    add_footer(s)

    # Slide 5
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s)
    add_title(s, "Plus-value SAP pour CashPilot")
    add_two_column(
        s,
        "Valeur operationnelle",
        [
            "Vision unique des priorites",
            "Execution guidee via workstreams",
            "Suivi factuel par score et statut",
            "Reduction des blocages et retards",
        ],
        "Valeur humaine",
        [
            "Lecture simplifiee pour novices",
            "Alignement CFO, compta, ops",
            "Moins de dependance a l'expert unique",
            "Montee en competence plus rapide",
        ],
    )
    add_footer(s)

    use_cases = [
        (
            1,
            "Demarrage d'une nouvelle societe",
            "La societe part de zero et doit fiabiliser sa base comptable rapidement.",
            "Prioriser FI dans le cockpit, creer workstreams sur plan comptable, mappings et ecritures initiales.",
            "Base comptable exploitable des le premier cycle de pilotage.",
            "Score FI, total done, reduction des workstreams planned.",
        ),
        (
            2,
            "Stabilisation des ecritures FI",
            "Incoherences entre pieces et ecritures, difficultes de suivi.",
            "Analyser metriques FI puis traiter les actions bloquees/overdue en roadmap.",
            "Qualite des ecritures amelioree et meilleure fiabilite du reporting.",
            "Entries, last entry, baisse des workstreams blocked.",
        ),
        (
            3,
            "Mise en place du controlling CO",
            "Absence d'axes analytiques pour suivre la performance.",
            "Activer CO, definir axes et deployer les regles analytiques par workstreams.",
            "Lecture des couts et marges par axe metier.",
            "Axes analytiques, score CO, taux de completion des actions.",
        ),
        (
            4,
            "Structuration des immobilisations AA",
            "Les actifs sont suivis en dehors du systeme.",
            "Lancer module AA, inventorier les actifs et planifier la reprise en base.",
            "Gestion des immobilisations centralisee et traçable.",
            "Assets count, score AA, workstreams done.",
        ),
        (
            5,
            "Preparation de cloture periodique",
            "Clotures tardives et charge de fin de mois trop elevee.",
            "Utiliser module Close pour orchestrer checklist et dependances.",
            "Cloture plus predicible et reduction des retards.",
            "Closures, latest closure, overdue proche de zero.",
        ),
        (
            6,
            "Preparation audit interne/externe",
            "Besoin de preuves de pilotage et de conformite.",
            "Presenter score global, statuts modules et roadmap execution.",
            "Narratif d'audit clair avec traces de progression.",
            "Global score, ratio done/blocked, historique generatedAt.",
        ),
        (
            7,
            "Pilotage multi-societes",
            "Croissance du groupe et complexite intercompany.",
            "Piloter module Consolidation et relier portfolios/members.",
            "Vision groupe harmonisee et meilleure maitrise des interco.",
            "Portfolios, members, score Consolidation.",
        ),
        (
            8,
            "Arbitrage capacite equipe finance",
            "Trop d'initiatives en parallele sans priorisation robuste.",
            "Trier workstreams par statut, priorite et overdue dans le cockpit.",
            "Roadmap realiste et sequence d'execution defendable.",
            "Blocked, overdue, completion_pct.",
        ),
        (
            9,
            "Onboarding d'un profil non-comptable",
            "Un nouveau collaborateur doit contribuer aux flux finance.",
            "Parcours guide via bulles info + modules + actions roadmap simples.",
            "Autonomie rapide sans surcharge de jargon technique.",
            "Temps d'onboarding, erreurs operationnelles, done par sprint.",
        ),
        (
            10,
            "Revue mensuelle de direction",
            "Direction souhaite un etat de progression finance lisible.",
            "Utiliser cockpit comme support de comite: score, risques, plan suivant.",
            "Decisions plus rapides et alignement transversal.",
            "Global score, top 3 risques, plan du mois suivant valide.",
        ),
    ]

    for uc in use_cases:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(s)
        add_use_case_slide(s, *uc)

    # Closing slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s)
    add_title(s, "Synthese")
    add_bullets(
        s,
        "Message cle",
        [
            "Le cockpit SAP de CashPilot transforme des modules comptables en execution guidee.",
            "Il facilite l'adoption des novices tout en restant utile aux experts.",
            "Il relie lecture KPI, priorisation roadmap et action operationnelle.",
            "Resultat: meilleure cadence de progression et meilleure qualite de pilotage.",
        ],
    )
    add_footer(s, "CashPilot • Merci")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Presentation generated: {OUTPUT}")


if __name__ == "__main__":
    build()
